#!/usr/bin/env python3
"""
文件文本提取工具

支持 docx、doc、pptx、pdf、txt、xlsx、xls、csv、tsv 格式
输出 JSON 格式的结构化文本内容

用法：
  python extract_file.py <文件路径>
  python extract_file.py <文件路径1> <文件路径2> ...
  python extract_file.py --dir <目录路径>

输出格式：
  [{"title": "文件名", "source": "文件路径", "content": "提取的文本", "pages": [...]}]

依赖安装：
  pip install python-docx python-pptx PyPDF2 openpyxl xlrd
"""

import sys
import os
import json
import csv


def extract_docx(filepath):
    """提取 docx 文件内容"""
    from docx import Document
    doc = Document(filepath)

    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    # 提取表格
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    return '\n'.join(paragraphs)


def extract_pptx(filepath):
    """提取 pptx 文件内容（逐页）"""
    from pptx import Presentation
    prs = Presentation(filepath)

    pages = []
    for i, slide in enumerate(prs.slides, 1):
        page_parts = []
        title = ''

        for shape in slide.shapes:
            # 标题
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    title = shape.text_frame.text.strip()
                else:
                    text = shape.text_frame.text.strip()
                    if text:
                        page_parts.append(text)

            # 表格
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        page_parts.append(row_text)

        # 备注
        notes = ''
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        page_content = ''
        if title:
            page_content += f'[第{i}页] {title}\n'
        if page_parts:
            page_content += '\n'.join(page_parts)
        if notes:
            page_content += f'\n备注：{notes}'

        if page_content.strip():
            pages.append({'page': i, 'title': title, 'content': page_content.strip()})

    full_text = '\n\n'.join(p['content'] for p in pages)
    return full_text, pages


def extract_pdf(filepath):
    """提取 PDF 文件内容"""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(filepath)
        pages = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                pages.append({'page': i, 'content': text.strip()})
        full_text = '\n\n'.join(p['content'] for p in pages)
        return full_text, pages
    except ImportError:
        # 备选方案
        try:
            import subprocess
            result = subprocess.run(['pdftotext', filepath, '-'], capture_output=True, text=True)
            if result.returncode == 0:
                return result.stdout, []
        except FileNotFoundError:
            pass
        return '', []


def extract_txt(filepath):
    """提取纯文本文件"""
    with open(filepath, 'r', encoding='utf-8') as f:
        return f.read()


def stringify_cell(value):
    """将表格单元格值转为适合审查的文本"""
    if value is None:
        return ''
    text = str(value).strip()
    return text


def extract_xlsx(filepath):
    """提取 xlsx/xlsm/xltx/xltm 文件内容（按工作表组织）"""
    from openpyxl import load_workbook

    workbook = load_workbook(filepath, read_only=True, data_only=True)
    sheets = []

    for sheet_index, worksheet in enumerate(workbook.worksheets, 1):
        rows = []
        for row_number, row in enumerate(worksheet.iter_rows(values_only=True), 1):
            values = [stringify_cell(cell) for cell in row]
            while values and not values[-1]:
                values.pop()
            if any(values):
                rows.append(f'第{row_number}行: ' + ' | '.join(values))

        if rows:
            content = f'[工作表{sheet_index}] {worksheet.title}\n' + '\n'.join(rows)
            sheets.append({'sheet': sheet_index, 'title': worksheet.title, 'content': content})

    full_text = '\n\n'.join(sheet['content'] for sheet in sheets)
    return full_text, sheets


def extract_xls(filepath):
    """提取 xls 文件内容（按工作表组织）"""
    import xlrd

    workbook = xlrd.open_workbook(filepath)
    sheets = []

    for sheet_index, worksheet in enumerate(workbook.sheets(), 1):
        rows = []
        for row_index in range(worksheet.nrows):
            values = [stringify_cell(worksheet.cell_value(row_index, col_index))
                      for col_index in range(worksheet.ncols)]
            while values and not values[-1]:
                values.pop()
            if any(values):
                rows.append(f'第{row_index + 1}行: ' + ' | '.join(values))

        if rows:
            content = f'[工作表{sheet_index}] {worksheet.name}\n' + '\n'.join(rows)
            sheets.append({'sheet': sheet_index, 'title': worksheet.name, 'content': content})

    full_text = '\n\n'.join(sheet['content'] for sheet in sheets)
    return full_text, sheets


def extract_delimited(filepath, delimiter):
    """提取 csv/tsv 文件内容"""
    encodings = ('utf-8-sig', 'utf-8', 'gb18030')
    last_error = None

    for encoding in encodings:
        try:
            rows = []
            with open(filepath, 'r', encoding=encoding, newline='') as f:
                reader = csv.reader(f, delimiter=delimiter)
                for row_number, row in enumerate(reader, 1):
                    values = [stringify_cell(cell) for cell in row]
                    while values and not values[-1]:
                        values.pop()
                    if any(values):
                        rows.append(f'第{row_number}行: ' + ' | '.join(values))
            return '\n'.join(rows)
        except UnicodeDecodeError as e:
            last_error = e

    raise last_error


def extract_file(filepath):
    """根据文件类型提取内容"""
    ext = os.path.splitext(filepath)[1].lower()
    filename = os.path.basename(filepath)

    result = {
        'title': os.path.splitext(filename)[0],
        'source': filepath,
        'content': '',
        'pages': [],
        'format': ext,
    }

    try:
        if ext in ('.docx',):
            result['content'] = extract_docx(filepath)
        elif ext in ('.pptx',):
            content, pages = extract_pptx(filepath)
            result['content'] = content
            result['pages'] = pages
        elif ext in ('.pdf',):
            content, pages = extract_pdf(filepath)
            result['content'] = content
            result['pages'] = pages
        elif ext in ('.txt', '.md'):
            result['content'] = extract_txt(filepath)
        elif ext in ('.xlsx', '.xlsm', '.xltx', '.xltm'):
            content, sheets = extract_xlsx(filepath)
            result['content'] = content
            result['pages'] = sheets
        elif ext in ('.xls',):
            content, sheets = extract_xls(filepath)
            result['content'] = content
            result['pages'] = sheets
        elif ext in ('.csv',):
            result['content'] = extract_delimited(filepath, ',')
        elif ext in ('.tsv',):
            result['content'] = extract_delimited(filepath, '\t')
        elif ext in ('.doc',):
            # .doc 格式尝试用 antiword 或 textutil (macOS)
            try:
                import subprocess
                r = subprocess.run(['textutil', '-convert', 'txt', '-stdout', filepath],
                                   capture_output=True, text=True)
                if r.returncode == 0:
                    result['content'] = r.stdout
            except Exception:
                result['content'] = ''
                result['error'] = '.doc 格式需要安装 antiword 或使用 macOS textutil'
        elif ext in ('.ppt',):
            result['content'] = ''
            result['error'] = '.ppt 格式建议转换为 .pptx 后处理'
        else:
            result['error'] = f'不支持的文件格式: {ext}'
    except Exception as e:
        result['error'] = str(e)

    return result


def scan_directory(dirpath):
    """递归扫描目录下所有支持的文件"""
    supported = (
        '.docx', '.pptx', '.pdf', '.txt', '.md', '.doc',
        '.xlsx', '.xlsm', '.xltx', '.xltm', '.xls', '.csv', '.tsv',
    )
    files = []
    for root, dirs, filenames in os.walk(dirpath):
        dirs[:] = [d for d in dirs if not d.startswith('.')]
        for f in sorted(filenames):
            if os.path.splitext(f)[1].lower() in supported and not f.startswith('~$'):
                files.append(os.path.join(root, f))
    return files


def main():
    args = sys.argv[1:]

    if not args:
        print('用法: python extract_file.py <文件路径> [文件路径2 ...]')
        print('      python extract_file.py --dir <目录路径>')
        sys.exit(0)

    files = []

    if args[0] == '--dir' and len(args) > 1:
        dirpath = args[1]
        files = scan_directory(dirpath)
        print(f'扫描到 {len(files)} 个文件', file=sys.stderr)
    else:
        files = [f for f in args if os.path.exists(f)]

    results = []
    for filepath in files:
        result = extract_file(filepath)
        if result['content']:
            results.append(result)
            print(f'✅ {os.path.basename(filepath)} ({len(result["content"])}字)', file=sys.stderr)
        else:
            print(f'⚠️  {os.path.basename(filepath)} (无内容或错误: {result.get("error", "")})', file=sys.stderr)

    # 输出 JSON 到 stdout
    print(json.dumps(results, ensure_ascii=False, indent=2))


if __name__ == '__main__':
    main()
